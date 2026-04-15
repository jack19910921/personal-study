#!/usr/bin/env python3
"""
PPT Maker — Generate beautiful PowerPoint presentations.

Usage:
    python create_pptx.py --topic "AI in Healthcare" --slides 8 --style corporate
    python create_pptx.py --outline outline.md --style startup
    python create_pptx.py --json slides.json --style minimal
    python create_pptx.py --list-styles
"""

import argparse
import json
import re
import sys
import os
from pathlib import Path
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.oxml.ns import qn

# ── Slide dimensions (16:9) ──
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ── Layout margins ──
MARGIN_LEFT = Inches(0.8)
MARGIN_RIGHT = Inches(0.8)
CONTENT_WIDTH = Inches(11.7)  # 13.333 - 0.8 - 0.8
TITLE_TOP = Inches(0.6)
BODY_TOP = Inches(1.3)

# ── Chart defaults ──
CHART_X = Inches(1.0)
CHART_Y = Inches(1.6)
CHART_W = Inches(11.3)
CHART_H = Inches(5.2)


# ============================================================
# STYLE PRESETS
# ============================================================
STYLES = {
    "corporate": {
        "name": "Corporate",
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title_color": RGBColor(0x1F, 0x3A, 0x5F),
        "body_color": RGBColor(0x33, 0x33, 0x33),
        "accent": RGBColor(0x4A, 0x90, 0xD9),
        "accent2": RGBColor(0xE7, 0x4C, 0x3C),
        "title_font": "Arial",
        "body_font": "Arial",
        "title_font_ea": "微软雅黑",
        "body_font_ea": "微软雅黑",
        "slide_bg": RGBColor(0xFF, 0xFF, 0xFF),
        "footer_color": RGBColor(0x99, 0x99, 0x99),
        "table_header": "1F3A5F",
        "table_alt_row": "F0F4F8",
        "dark_slide_bg": RGBColor(0x1F, 0x3A, 0x5F),
        "dark_text": RGBColor(0xFF, 0xFF, 0xFF),
        "dark_subtext": RGBColor(0xCC, 0xDD, 0xEE),
        "sizes": {"title": 44, "subtitle": 20, "heading": 30, "body": 18, "caption": 12},
    },
    "startup": {
        "name": "Startup Pitch",
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title_color": RGBColor(0x2D, 0x34, 0x36),
        "body_color": RGBColor(0x63, 0x6E, 0x72),
        "accent": RGBColor(0x6C, 0x5C, 0xE7),
        "accent2": RGBColor(0x00, 0xB8, 0x94),
        "title_font": "Segoe UI",
        "body_font": "Segoe UI",
        "title_font_ea": "微软雅黑",
        "body_font_ea": "微软雅黑",
        "slide_bg": RGBColor(0xFF, 0xFF, 0xFF),
        "footer_color": RGBColor(0xAA, 0xAA, 0xAA),
        "table_header": "6C5CE7",
        "table_alt_row": "F3F0FF",
        "dark_slide_bg": RGBColor(0x2D, 0x34, 0x36),
        "dark_text": RGBColor(0xFF, 0xFF, 0xFF),
        "dark_subtext": RGBColor(0xDD, 0xDD, 0xEE),
        "sizes": {"title": 48, "subtitle": 22, "heading": 32, "body": 18, "caption": 12},
    },
    "executive": {
        "name": "Executive",
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title_color": RGBColor(0x1E, 0x3A, 0x5F),
        "body_color": RGBColor(0x2C, 0x3E, 0x50),
        "accent": RGBColor(0xC9, 0xA2, 0x27),
        "accent2": RGBColor(0x1E, 0x3A, 0x5F),
        "title_font": "Georgia",
        "body_font": "Calibri",
        "title_font_ea": "微软雅黑",
        "body_font_ea": "微软雅黑",
        "slide_bg": RGBColor(0xFA, 0xF9, 0xF6),
        "footer_color": RGBColor(0x99, 0x99, 0x99),
        "table_header": "1E3A5F",
        "table_alt_row": "F5F2EB",
        "dark_slide_bg": RGBColor(0x1E, 0x3A, 0x5F),
        "dark_text": RGBColor(0xFF, 0xFF, 0xFF),
        "dark_subtext": RGBColor(0xDD, 0xCC, 0x99),
        "sizes": {"title": 42, "subtitle": 22, "heading": 28, "body": 18, "caption": 12},
    },
    "dark": {
        "name": "Dark Mode",
        "bg": RGBColor(0x1A, 0x1A, 0x2E),
        "title_color": RGBColor(0xFF, 0xFF, 0xFF),
        "body_color": RGBColor(0xE0, 0xE0, 0xE0),
        "accent": RGBColor(0x00, 0xD4, 0xFF),
        "accent2": RGBColor(0xFF, 0x6B, 0x6B),
        "title_font": "Segoe UI",
        "body_font": "Segoe UI",
        "title_font_ea": "微软雅黑",
        "body_font_ea": "微软雅黑",
        "slide_bg": RGBColor(0x1A, 0x1A, 0x2E),
        "footer_color": RGBColor(0x88, 0x88, 0x88),
        "table_header": "00D4FF",
        "table_alt_row": "252540",
        "dark_slide_bg": RGBColor(0x12, 0x12, 0x22),
        "dark_text": RGBColor(0xFF, 0xFF, 0xFF),
        "dark_subtext": RGBColor(0xBB, 0xBB, 0xBB),
        "sizes": {"title": 44, "subtitle": 22, "heading": 30, "body": 18, "caption": 12},
    },
    "minimal": {
        "name": "Minimal",
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title_color": RGBColor(0x1A, 0x1A, 0x1A),
        "body_color": RGBColor(0x4A, 0x4A, 0x4A),
        "accent": RGBColor(0x00, 0x66, 0xCC),
        "accent2": RGBColor(0x00, 0xA8, 0x6B),
        "title_font": "Segoe UI",
        "body_font": "Segoe UI",
        "title_font_ea": "微软雅黑",
        "body_font_ea": "微软雅黑",
        "slide_bg": RGBColor(0xFF, 0xFF, 0xFF),
        "footer_color": RGBColor(0x99, 0x99, 0x99),
        "table_header": "0066CC",
        "table_alt_row": "F5F8FC",
        "dark_slide_bg": RGBColor(0x1A, 0x1A, 0x1A),
        "dark_text": RGBColor(0xFF, 0xFF, 0xFF),
        "dark_subtext": RGBColor(0xCC, 0xCC, 0xCC),
        "sizes": {"title": 44, "subtitle": 20, "heading": 30, "body": 18, "caption": 12},
    },
    "creative": {
        "name": "Creative",
        "bg": RGBColor(0xFA, 0xFA, 0xFA),
        "title_color": RGBColor(0x2D, 0x2D, 0x2D),
        "body_color": RGBColor(0x4A, 0x4A, 0x4A),
        "accent": RGBColor(0xFF, 0x57, 0x22),
        "accent2": RGBColor(0x9C, 0x27, 0xB0),
        "title_font": "Segoe UI",
        "body_font": "Segoe UI",
        "title_font_ea": "微软雅黑",
        "body_font_ea": "微软雅黑",
        "slide_bg": RGBColor(0xFA, 0xFA, 0xFA),
        "footer_color": RGBColor(0x99, 0x99, 0x99),
        "table_header": "FF5722",
        "table_alt_row": "FFF3E0",
        "dark_slide_bg": RGBColor(0x2D, 0x2D, 0x2D),
        "dark_text": RGBColor(0xFF, 0xFF, 0xFF),
        "dark_subtext": RGBColor(0xDD, 0xDD, 0xDD),
        "sizes": {"title": 48, "subtitle": 22, "heading": 32, "body": 20, "caption": 12},
    },
}


# ============================================================
# CHART TYPE MAP
# ============================================================
CHART_MAP = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar_stacked": XL_CHART_TYPE.COLUMN_STACKED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "column_stacked": XL_CHART_TYPE.COLUMN_STACKED,
    "line": XL_CHART_TYPE.LINE,
    "line_markers": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
    "area": XL_CHART_TYPE.AREA_STACKED,
    "area_stacked": XL_CHART_TYPE.AREA_STACKED,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
}


# ============================================================
# UTILITY FUNCTIONS
# ============================================================
def set_cell_shading(cell, color_hex: str):
    """Set background color for a table cell."""
    tc = cell._tc
    tcPr = tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = tc.makeelement(qn('a:tcPr'), {})
        tc.insert(0, tcPr)
    solidFill = tcPr.makeelement(qn('a:solidFill'), {})
    srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': color_hex})
    solidFill.append(srgbClr)
    tcPr.append(solidFill)


def set_font(run_or_para, style, bold=False, italic=False, size=None):
    """Apply font with Chinese font support."""
    if size:
        if hasattr(run_or_para, 'font'):
            run_or_para.font.size = Pt(size)
    if bold:
        run_or_para.font.bold = True
    if italic:
        run_or_para.font.italic = True

    # Chinese font via XML
    if hasattr(run_or_para, '_r'):
        r = run_or_para._r
    elif hasattr(run_or_para, 'runs') and run_or_para.runs:
        r = run_or_para.runs[0]._r
    else:
        return

    rPr = r.get_or_add_rPr()
    ea = rPr.makeelement(qn('a:eaFont'), {'typeface': style.get('body_font_ea', '微软雅黑')})
    rPr.append(ea)
    latin = rPr.makeelement(qn('a:latin'), {'typeface': style.get('body_font', 'Segoe UI')})
    rPr.append(latin)


def apply_para_style(para, style, size_key='body', color_key='body_color'):
    """Apply consistent paragraph styling."""
    sz = style['sizes'][size_key]
    color = style[color_key]
    para.font.size = Pt(sz)
    para.font.color.rgb = color
    para.font.name = style['body_font']
    para.space_after = Pt(6)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, style, top=None, height=Inches(0.12), full_width=True):
    """Add accent color bar at top or bottom of slide."""
    w = SLIDE_WIDTH if full_width else Inches(3.0)
    if top is None:
        top = SLIDE_HEIGHT - height
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), top, w, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = style['accent']
    bar.line.fill.background()


def add_footer(slide, style, slide_num, total_slides):
    """Add slide number footer."""
    txBox = slide.shapes.add_textbox(
        Inches(0.8), SLIDE_HEIGHT - Inches(0.45),
        Inches(5), Inches(0.3)
    )
    tf = txBox.text_frame
    tf.text = f"{slide_num} / {total_slides}"
    tf.paragraphs[0].font.size = Pt(style['sizes']['caption'])
    tf.paragraphs[0].font.color.rgb = style['footer_color']
    tf.paragraphs[0].font.name = style['body_font']

    # Date on right
    txBox2 = slide.shapes.add_textbox(
        Inches(8), SLIDE_HEIGHT - Inches(0.45),
        Inches(4.5), Inches(0.3)
    )
    tf2 = txBox2.text_frame
    tf2.text = datetime.now().strftime("%Y-%m-%d")
    tf2.paragraphs[0].font.size = Pt(style['sizes']['caption'])
    tf2.paragraphs[0].font.color.rgb = style['footer_color']
    tf2.paragraphs[0].font.name = style['body_font']
    tf2.paragraphs[0].alignment = PP_ALIGN.RIGHT


def count_slides_data(data):
    """Count total slides for footer."""
    return len(data.get('slides', [])) + (1 if data.get('title') else 0) + (1 if data.get('ending') != False else 0)


# ============================================================
# SLIDE BUILDERS
# ============================================================
def build_title_slide(prs, style, data):
    """Build cover/title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, style['dark_slide_bg'])

    title = slide.shapes.title
    if title:
        title.text = data['title']
        for p in title.text_frame.paragraphs:
            p.font.size = Pt(style['sizes']['title'])
            p.font.bold = True
            p.font.color.rgb = style['dark_text']
            p.font.name = style['title_font']
            p.alignment = PP_ALIGN.CENTER

    subtitle = slide.placeholders[1] if len(slide.placeholders) > 1 else None
    if subtitle and data.get('subtitle') is not None and data['subtitle']:
        sub_text = data['subtitle']
        if data.get('author'):
            sub_text += f"\n{data['author']}"
        subtitle.text = sub_text
        for p in subtitle.text_frame.paragraphs:
            p.font.size = Pt(style['sizes']['subtitle'])
            p.font.color.rgb = style['dark_subtext']
            p.font.name = style['body_font']
            p.alignment = PP_ALIGN.CENTER

    # Top accent bar
    add_accent_bar(slide, style, top=Inches(0), height=Inches(0.12))


def build_bullet_slide(prs, style, slide_data, slide_num, total):
    """Build title + bullet points slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_bg(slide, style['slide_bg'])

    if slide.shapes.title and slide_data.get('title'):
        slide.shapes.title.text = slide_data['title']
        for p in slide.shapes.title.text_frame.paragraphs:
            p.font.size = Pt(style['sizes']['heading'])
            p.font.bold = True
            p.font.color.rgb = style['title_color']
            p.font.name = style['title_font']

    # Accent bar under title
    add_accent_bar(slide, style, top=Inches(1.05), height=Inches(0.06), full_width=False)

    # Bullets
    bullets = slide_data.get('bullets', [])
    if bullets:
        content = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                content = shape
                break

        if content and hasattr(content, 'text_frame'):
            tf = content.text_frame
            tf.clear()

            for i, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = bullet
                apply_para_style(p, style)
                p.level = bullet.get('level', 0) if isinstance(bullet, dict) else 0
                # Bold the first bullet if it looks like a sub-heading
                if i == 0 and len(bullets) > 3:
                    p.font.bold = True
                    p.font.size = Pt(style['sizes']['body'] + 2)

    add_footer(slide, style, slide_num, total)


def build_section_slide(prs, style, slide_data, slide_num, total):
    """Build section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, style['dark_slide_bg'])

    txBox = slide.shapes.add_textbox(
        Inches(1.5), Inches(2.5), Inches(10.3), Inches(2.5)
    )
    tf = txBox.text_frame
    if slide_data.get('title'):
        tf.text = slide_data['title']
        tf.paragraphs[0].font.size = Pt(style['sizes']['title'] + 4)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = style['accent']
        tf.paragraphs[0].font.name = style['title_font']
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    if slide_data.get('subtitle'):
        p = tf.add_paragraph()
        p.text = slide_data['subtitle']
        p.font.size = Pt(style['sizes']['subtitle'])
        p.font.color.rgb = style['dark_subtext']
        p.font.name = style['body_font']
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(20)

    add_accent_bar(slide, style, top=Inches(0))


def build_two_column_slide(prs, style, slide_data, slide_num, total):
    """Build two-column layout slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    set_slide_bg(slide, style['slide_bg'])

    if slide.shapes.title and slide_data.get('title'):
        slide.shapes.title.text = slide_data['title']
        for p in slide.shapes.title.text_frame.paragraphs:
            p.font.size = Pt(style['sizes']['heading'])
            p.font.bold = True
            p.font.color.rgb = style['title_color']
            p.font.name = style['title_font']

    add_accent_bar(slide, style, top=Inches(1.05), height=Inches(0.06), full_width=False)

    col_w = 5.4
    gap = 0.8
    bullets = slide_data.get('bullets', [])
    mid = len(bullets) // 2

    for col_idx, (left_pos, items) in enumerate([
        (Inches(0.8), bullets[:mid]),
        (Inches(0.8 + col_w + gap), bullets[mid:]),
    ]):
        if not items:
            continue
        txBox = slide.shapes.add_textbox(left_pos, Inches(1.4), Inches(col_w), Inches(5.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            apply_para_style(p, style)
            p.space_before = Pt(8)

    add_footer(slide, style, slide_num, total)


def build_image_text_slide(prs, style, slide_data, slide_num, total):
    """Build image + text layout slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_bg(slide, style['slide_bg'])

    if slide.shapes.title and slide_data.get('title'):
        slide.shapes.title.text = slide_data['title']
        for p in slide.shapes.title.text_frame.paragraphs:
            p.font.size = Pt(style['sizes']['heading'])
            p.font.bold = True
            p.font.color.rgb = style['title_color']
            p.font.name = style['title_font']

    add_accent_bar(slide, style, top=Inches(1.05), height=Inches(0.06), full_width=False)

    # Image placeholder (left side)
    img_path = slide_data.get('image_path')
    if img_path and Path(img_path).exists():
        slide.shapes.add_picture(img_path, Inches(0.8), Inches(1.5), width=Inches(6.0))
    else:
        placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(6.0), Inches(4.5)
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = style['accent']
        placeholder.fill.fore_color.brightness = 0.85
        placeholder.line.color.rgb = style['accent']
        placeholder.text_frame.paragraphs[0].text = "[ Image Placeholder ]"
        placeholder.text_frame.paragraphs[0].font.color.rgb = style['title_color']
        placeholder.text_frame.paragraphs[0].font.size = Pt(16)
        placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Text (right side)
    bullets = slide_data.get('bullets', [])
    txBox = slide.shapes.add_textbox(Inches(7.2), Inches(1.5), Inches(5.3), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        apply_para_style(p, style)
        p.space_before = Pt(10)

    add_footer(slide, style, slide_num, total)


def build_table_slide(prs, style, slide_data, slide_num, total):
    """Build data table slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_bg(slide, style['slide_bg'])

    if slide.shapes.title and slide_data.get('title'):
        slide.shapes.title.text = slide_data['title']
        for p in slide.shapes.title.text_frame.paragraphs:
            p.font.size = Pt(style['sizes']['heading'])
            p.font.bold = True
            p.font.color.rgb = style['title_color']
            p.font.name = style['title_font']

    add_accent_bar(slide, style, top=Inches(1.05), height=Inches(0.06), full_width=False)

    # Table
    table_data = slide_data.get('table', {})
    headers = table_data.get('headers', [])
    rows = table_data.get('rows', [])

    if not headers:
        if rows:
            headers = [f"Col {i+1}" for i in range(len(rows[0]))]
        else:
            headers = ["Empty"]
            rows = [["No data available"]]

    n_rows = len(rows) + 1
    n_cols = len(headers)

    shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.5 * n_rows + 0.5)
    )
    table = shape.table

    # Header
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        set_cell_shading(cell, style['table_header'])

    # Rows
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = style['body_color']
            if r_idx % 2 == 0:
                set_cell_shading(cell, style['table_alt_row'])

    add_footer(slide, style, slide_num, total)


def build_chart_slide(prs, style, slide_data, slide_num, total):
    """Build chart slide — actually generates real charts."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_bg(slide, style['slide_bg'])

    if slide.shapes.title and slide_data.get('title'):
        slide.shapes.title.text = slide_data['title']
        for p in slide.shapes.title.text_frame.paragraphs:
            p.font.size = Pt(style['sizes']['heading'])
            p.font.bold = True
            p.font.color.rgb = style['title_color']
            p.font.name = style['title_font']

    add_accent_bar(slide, style, top=Inches(1.05), height=Inches(0.06), full_width=False)

    chart_data_def = slide_data.get('chart_data', {})
    chart_type = slide_data.get('chart_type', 'bar')

    # Scatter charts need XyChartData
    is_scatter = chart_type in ('scatter',)
    if is_scatter:
        chart_data = XyChartData()
        series_list = chart_data_def.get('series', [])
        if series_list:
            for s in series_list:
                series = chart_data.add_series(s.get('name', 'Series'))
                for pt in s.get('points', []):
                    if isinstance(pt, (list, tuple)) and len(pt) >= 2:
                        series.add_data_point(pt[0], pt[1])
                    elif isinstance(pt, dict):
                        series.add_data_point(pt.get('x', 0), pt.get('y', 0))
        else:
            # Fallback scatter data
            series = chart_data.add_series('Data')
            for x, y in [(1, 2), (2, 4), (3, 3), (4, 5), (5, 4)]:
                series.add_data_point(x, y)
    else:
        chart_data = CategoryChartData()
        categories = chart_data_def.get('categories', [])
        series_list = chart_data_def.get('series', [])

        if categories and series_list:
            chart_data.categories = categories
            for s in series_list:
                chart_data.add_series(s.get('name', 'Series'), tuple(s.get('values', [])))
        elif chart_data_def.get('data'):
            # Inline data: {"Q1": 120, "Q2": 145, ...}
            chart_data.categories = list(chart_data_def['data'].keys())
            chart_data.add_series('Values', tuple(chart_data_def['data'].values()))

        if not chart_data.categories:
            # Fallback data
            chart_data.categories = ['A', 'B', 'C', 'D']
            chart_data.add_series('Data', (10, 20, 15, 25))

    chart_type_enum = CHART_MAP.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

    gf = slide.shapes.add_chart(chart_type_enum, CHART_X, CHART_Y, CHART_W, CHART_H, chart_data)
    chart = gf.chart

    # Legend
    is_pie = chart_type in ('pie', 'doughnut')
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM if not is_pie else XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(13)
    chart.legend.font.color.rgb = style['body_color']

    # Data labels (scatter charts don't support this via simple API)
    if not is_pie and not is_scatter:
        plot = chart.plots[0]
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.font.size = Pt(11)
        dl.font.color.rgb = style['body_color']

    # Axes (category axis only for non-scatter, non-pie charts)
    if not is_pie and not is_scatter:
        chart.category_axis.tick_labels.font.size = Pt(12)
        chart.category_axis.tick_labels.font.color.rgb = style['body_color']
        chart.value_axis.tick_labels.font.size = Pt(11)
        chart.value_axis.tick_labels.font.color.rgb = style['body_color']

    add_footer(slide, style, slide_num, total)


def build_content_slide(prs, style, slide_data, slide_num, total):
    """Build title + free-form content slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_bg(slide, style['slide_bg'])

    if slide.shapes.title and slide_data.get('title'):
        slide.shapes.title.text = slide_data['title']
        for p in slide.shapes.title.text_frame.paragraphs:
            p.font.size = Pt(style['sizes']['heading'])
            p.font.bold = True
            p.font.color.rgb = style['title_color']
            p.font.name = style['title_font']

    add_accent_bar(slide, style, top=Inches(1.05), height=Inches(0.06), full_width=False)

    # Content box
    txBox = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    content_lines = slide_data.get('content', [])
    for i, line in enumerate(content_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        apply_para_style(p, style)
        p.space_before = Pt(12)

    add_footer(slide, style, slide_num, total)


def build_ending_slide(prs, style, slide_data=None, slide_num=None, total=None):
    """Build Thank You / ending slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, style['dark_slide_bg'])

    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.3), Inches(2.5))
    tf = txBox.text_frame
    tf.text = "Thank You"
    tf.paragraphs[0].font.size = Pt(style['sizes']['title'] + 10)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = style['accent']
    tf.paragraphs[0].font.name = style['title_font']
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Questions?"
    p.font.size = Pt(style['sizes']['subtitle'])
    p.font.color.rgb = style['dark_subtext']
    p.font.name = style['body_font']
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(20)

    if slide_data and slide_data.get('contact'):
        p2 = tf.add_paragraph()
        p2.text = slide_data['contact']
        p2.font.size = Pt(style['sizes']['body'])
        p2.font.color.rgb = style['footer_color']
        p2.font.name = style['body_font']
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(16)

    add_accent_bar(slide, style, top=Inches(0))


# ============================================================
# OUTLINE PARSER
# ============================================================
def parse_outline(outline_path: str) -> dict:
    """Parse markdown outline into slide structure."""
    with open(outline_path, 'r') as f:
        content = f.read()

    presentation = {
        "title": "",
        "subtitle": "",
        "author": "",
        "slides": []
    }

    lines = content.strip().split("\n")
    current_slide = None

    for line in lines:
        line = line.rstrip()

        if line.startswith("# ") and not presentation["title"]:
            presentation["title"] = line[2:].strip()
            continue

        if line.lower().startswith("subtitle:"):
            presentation["subtitle"] = line.split(":", 1)[1].strip()
            continue
        if line.lower().startswith("author:"):
            presentation["author"] = line.split(":", 1)[1].strip()
            continue

        if line.startswith("## "):
            if current_slide:
                presentation["slides"].append(current_slide)

            title = line[3:].strip()
            title = re.sub(r"^Slide\s+\d+:\s*", "", title, flags=re.IGNORECASE)

            current_slide = {
                "title": title,
                "layout": "title_and_content",
                "bullets": [],
                "notes": "",
                "image_path": None,
                "chart_type": None,
                "chart_data": None,
                "table": None,
            }
            continue

        if current_slide:
            if line.startswith("- "):
                item = line[2:].strip()

                if item.lower().startswith("chart:") or item.lower().startswith("layout: chart"):
                    chart_type = item.split(":", 1)[1].strip() if ":" in item else "bar"
                    current_slide["chart_type"] = chart_type
                    current_slide["layout"] = "chart"
                elif item.lower().startswith("table:"):
                    current_slide["table"] = {"name": item.split(":", 1)[1].strip()}
                    current_slide["layout"] = "table"
                elif item.lower().startswith("data:"):
                    # Inline data: data: Q1=120, Q2=145, Q3=132
                    data_str = item.split(":", 1)[1].strip()
                    parsed = {}
                    for pair in data_str.split(","):
                        if "=" in pair:
                            k, v = pair.split("=", 1)
                            try:
                                parsed[k.strip()] = float(v.strip())
                            except ValueError:
                                parsed[k.strip()] = v.strip()
                    current_slide["chart_data"] = {"data": parsed}
                    if not current_slide["chart_type"]:
                        current_slide["chart_type"] = "bar"
                        current_slide["layout"] = "chart"
                elif item.startswith("!["):
                    match = re.match(r"!\[([^\]]*)\]\(([^)]+)\)", item)
                    if match:
                        alt, src = match.groups()
                        current_slide["image_path"] = src.strip()
                        current_slide["layout"] = "image_and_text"
                elif item.startswith("layout:"):
                    current_slide["layout"] = item.split(":", 1)[1].strip()
                elif item.startswith("image:"):
                    current_slide["image_path"] = item.split(":", 1)[1].strip()
                    current_slide["layout"] = "image_and_text"
                else:
                    current_slide["bullets"].append(item)

            elif line.strip().startswith(">"):
                current_slide["notes"] += line.strip()[1:].strip() + "\n"

    if current_slide:
        presentation["slides"].append(current_slide)

    return presentation


# ============================================================
# MAIN CREATION ENGINE
# ============================================================
def create_presentation(data: dict, style_name: str = "corporate") -> Presentation:
    """Create PowerPoint presentation from structured data."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    style = STYLES.get(style_name, STYLES["corporate"])

    # Count total slides for footer
    total = count_slides_data(data)

    slide_num = 0

    # 1. Title slide
    if data.get("title"):
        slide_num += 1
        build_title_slide(prs, style, data)

    # 2. Content slides
    for slide_data in data.get("slides", []):
        slide_num += 1
        layout = slide_data.get("layout", "title_and_content")

        if layout == "section":
            build_section_slide(prs, style, slide_data, slide_num, total)
        elif layout == "two_column":
            build_two_column_slide(prs, style, slide_data, slide_num, total)
        elif layout == "image_and_text":
            build_image_text_slide(prs, style, slide_data, slide_num, total)
        elif layout == "table":
            build_table_slide(prs, style, slide_data, slide_num, total)
        elif layout == "chart":
            build_chart_slide(prs, style, slide_data, slide_num, total)
        elif layout == "blank":
            # Build a content slide
            build_content_slide(prs, style, slide_data, slide_num, total)
        else:
            # Default: bullet slide
            build_bullet_slide(prs, style, slide_data, slide_num, total)

    # 3. Ending slide
    if data.get("ending", True):
        slide_num += 1
        ending_data = data.get("ending_data", {})
        build_ending_slide(prs, style, ending_data)

    return prs


# ============================================================
# AUTO OUTLINE GENERATOR
# ============================================================
def generate_outline_from_topic(topic: str, n_slides: int = 6) -> dict:
    """Generate a sensible outline from a topic string."""
    sections = [
        {"title": "Overview", "layout": "title_and_content", "bullets": [
            "What is this about?", "Why it matters now", "Key objectives"]},
        {"title": "Background", "layout": "title_and_content", "bullets": [
            "Current landscape", "Key challenges", "Opportunity areas"]},
        {"title": "Analysis", "layout": "title_and_content", "bullets": [
            "Data-driven insights", "Key findings", "Implications"]},
        {"title": "Solutions", "layout": "title_and_content", "bullets": [
            "Recommended approach", "Implementation roadmap", "Expected outcomes"]},
        {"title": "Summary", "layout": "section", "bullets": [
            "Key takeaways", "Next steps"]},
    ]

    return {
        "title": topic,
        "subtitle": f"Deep Dive & Analysis",
        "slides": sections[:n_slides - 2],  # Reserve for title + ending
    }


# ============================================================
# CLI
# ============================================================
def main():
    parser = argparse.ArgumentParser(description="Create beautiful PowerPoint presentations")
    parser.add_argument("--outline", "-o", help="Markdown outline file")
    parser.add_argument("--json", "-j", help="JSON structure file")
    parser.add_argument("--topic", "-t", help="Topic for auto-generated outline")
    parser.add_argument("--slides", "-n", type=int, default=6, help="Number of slides")
    parser.add_argument("--style", "-s", default="corporate", help=f"Style preset (default: corporate)")
    parser.add_argument("--list-styles", action="store_true", help="List available styles")
    parser.add_argument("--output", "-O", default="presentation.pptx", help="Output file path")

    args = parser.parse_args()

    if args.list_styles:
        print("Available styles:")
        for name, info in STYLES.items():
            print(f"  {name:15s} — {info['name']}")
        return

    data = None
    if args.outline:
        data = parse_outline(args.outline)
    elif args.json:
        with open(args.json) as f:
            data = json.load(f)
    elif args.topic:
        data = generate_outline_from_topic(args.topic, args.slides)
        print(f"Note: Generated outline for '{args.topic}' — edit the source for real content.", file=sys.stderr)
    else:
        print("Error: Provide --outline, --json, or --topic", file=sys.stderr)
        sys.exit(1)

    prs = create_presentation(data, style_name=args.style)
    prs.save(args.output)
    print(f"Created: {args.output}  ({len(prs.slides)} slides, style: {args.style})")


if __name__ == "__main__":
    main()
