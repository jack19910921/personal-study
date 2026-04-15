#!/usr/bin/env python3
"""Analyze a PowerPoint file to extract layouts, fonts, colors, and structure."""

import argparse
import json
from pathlib import Path

from pptx import Presentation


def analyze_pptx(path: str, verbose: bool = False) -> dict:
    prs = Presentation(path)
    result = {
        "file": str(path),
        "dimensions": {
            "width_inches": round(prs.slide_width.inches, 2),
            "height_inches": round(prs.slide_height.inches, 2),
        },
        "slide_count": len(prs.slides),
        "masters": [],
        "layouts": [],
        "slides": [],
    }

    for i, master in enumerate(prs.slide_masters):
        master_info = {"index": i, "layout_count": len(master.slide_layouts), "layouts": []}
        for j, layout in enumerate(master.slide_layouts):
            layout_info = {"index": j, "name": layout.name, "placeholders": []}
            for ph in layout.placeholders:
                layout_info["placeholders"].append({
                    "idx": ph.placeholder_format.idx,
                    "type": str(ph.placeholder_format.type),
                    "name": ph.name,
                })
            master_info["layouts"].append(layout_info)
            result["layouts"].append({"name": layout.name, "index": j, "master": i})
        result["masters"].append(master_info)

    max_slides = len(prs.slides) if verbose else min(10, len(prs.slides))
    for idx in range(max_slides):
        slide = prs.slides[idx]
        slide_info = {"number": idx + 1, "layout": slide.slide_layout.name, "shapes": []}
        for shape in slide.shapes:
            s = {"name": shape.name, "type": str(shape.shape_type)}
            if shape.left:
                s["left"] = round(shape.left.inches, 2)
            if shape.top:
                s["top"] = round(shape.top.inches, 2)
            if shape.width:
                s["width"] = round(shape.width.inches, 2)
            if shape.height:
                s["height"] = round(shape.height.inches, 2)
            if hasattr(shape, 'text') and shape.text.strip():
                s["text"] = shape.text.strip()[:100]
            if hasattr(shape, 'fill'):
                try:
                    if shape.fill.type is not None:
                        s["fill_type"] = str(shape.fill.type)
                        rgb = shape.fill.fore_color.rgb
                        if rgb:
                            s["fill_color"] = f"#{rgb}"
                except:
                    pass
            slide_info["shapes"].append(s)
        result["slides"].append(slide_info)

    return result


def print_summary(r: dict):
    print(f"PPTX Analysis: {Path(r['file']).name}")
    print(f"  Dimensions: {r['dimensions']['width_inches']}\" x {r['dimensions']['height_inches']}\"")
    print(f"  Total slides: {r['slide_count']}")
    print(f"\n  Available Layouts:")
    for l in r["layouts"]:
        print(f"    [{l['index']}] {l['name']}")
    print(f"\n  Slides (first {len(r['slides'])}):")
    for s in r["slides"][:10]:
        print(f"\n    Slide {s['number']}: {s['layout']}")
        for sh in s["shapes"][:5]:
            txt = f" -> \"{sh.get('text', '')[:40]}\"" if sh.get('text') else ""
            print(f"      - {sh['type']}: {sh['name']}{txt}")


def main():
    parser = argparse.ArgumentParser(description="Analyze PowerPoint file")
    parser.add_argument("file", help="Path to PPTX file")
    parser.add_argument("--json", "-j", action="store_true", help="Output as JSON")
    parser.add_argument("--verbose", "-v", action="store_true", help="Analyze all slides")
    parser.add_argument("--output", "-o", help="Save JSON to file")
    args = parser.parse_args()

    result = analyze_pptx(args.file, verbose=args.verbose)

    if args.json or args.output:
        out = json.dumps(result, indent=2)
        if args.output:
            Path(args.output).write_text(out)
            print(f"Saved to {args.output}")
        else:
            print(out)
    else:
        print_summary(result)


if __name__ == "__main__":
    main()
